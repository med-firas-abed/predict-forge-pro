from __future__ import annotations

import importlib.util
import json
import subprocess
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
GUIDES = Path(__file__).resolve().parents[1]
FINAL_REPORT = ROOT / "final_report"
IMG = FINAL_REPORT / "prediteq_overleaf_report" / "images"
IMG_FREE = FINAL_REPORT / "prediteq_overleaf_report" / "images_free"
OUTPUT_PPTX = GUIDES / "slides" / "PrediTeq_Soutenance_Jury_Generated.pptx"
OUTPUT_PDF = GUIDES / "slides" / "PrediTeq_Soutenance_Jury_Generated.pdf"


NAVY = RGBColor(11, 19, 43)
STEEL = RGBColor(28, 37, 65)
SLATE = RGBColor(71, 85, 105)
INK = RGBColor(30, 41, 59)
PAPER = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
ORANGE = RGBColor(243, 114, 44)
TEAL = RGBColor(42, 157, 143)
SKY = RGBColor(59, 130, 246)
GOLD = RGBColor(234, 179, 8)
MUTED = RGBColor(100, 116, 139)
GREEN = RGBColor(22, 163, 74)
RED = RGBColor(220, 38, 38)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def load_module(path: Path, module_name: str):
    spec = importlib.util.spec_from_file_location(module_name, path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Cannot load module from {path}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


CFG = load_module(ROOT / "prediteq_ml" / "config.py", "prediteq_config")
DEMO_CONTEXT = load_module(ROOT / "prediteq_api" / "demo_scenarios.py", "prediteq_demo_scenarios")
DEMO_MACHINE_SCENARIOS = DEMO_CONTEXT.DEMO_MACHINE_SCENARIOS
METRICS = load_json(ROOT / "prediteq_ml" / "outputs" / "metrics.json")
RUL_CV = load_json(ROOT / "prediteq_ml" / "outputs" / "rul_cv_scores.json")
CMAPSS = load_json(ROOT / "prediteq_ml" / "outputs" / "cmapss_metrics.json")
LEAD = load_json(ROOT / "prediteq_ml" / "outputs" / "lead_time.json")


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_name="Aptos",
    font_size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets,
    *,
    font_name="Aptos",
    font_size=19,
    color=INK,
    bullet_color=ORANGE,
    leading_text=None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if leading_text:
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.LEFT
        run = p0.add_run()
        run.text = leading_text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = color
    else:
        tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.add_paragraph() if leading_text or idx > 0 else tf.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(5)
    return box


def add_title_block(slide, title: str, subtitle: str | None = None, dark=False):
    color = WHITE if dark else INK
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.45), Inches(0.12), Inches(0.72))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.color.rgb = ORANGE
    add_text(slide, Inches(0.88), Inches(0.34), Inches(10.8), Inches(0.55), title, font_name="Aptos Display", font_size=28, color=color, bold=True)
    if subtitle:
        add_text(slide, Inches(0.9), Inches(0.85), Inches(11.2), Inches(0.34), subtitle, font_size=12, color=MUTED if not dark else RGBColor(203, 213, 225))


def add_footer(slide, page_num: int, total: int):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(11.5), Inches(0.25), "PrediTeq | Soutenance jury | simulation -> ML -> runtime -> web app", font_size=9, color=MUTED if slide.background.fill.fore_color.rgb != NAVY else RGBColor(203, 213, 225))
    add_text(slide, Inches(12.0), Inches(7.02), Inches(0.7), Inches(0.25), f"{page_num}/{total}", font_size=10, color=MUTED if slide.background.fill.fore_color.rgb != NAVY else RGBColor(203, 213, 225), align=PP_ALIGN.RIGHT)


def add_tag(slide, left, top, text, fill_color, text_color=WHITE):
    shape = add_rect(slide, left, top, Inches(1.28), Inches(0.34), fill_color)
    add_text(slide, left, top + Inches(0.02), Inches(1.28), Inches(0.22), text, font_size=10, color=text_color, bold=True, align=PP_ALIGN.CENTER)
    return shape


def add_stat_card(slide, left, top, width, height, value, label, accent, dark=False):
    fill_color = STEEL if dark else WHITE
    text_color = WHITE if dark else INK
    border = accent
    card = add_rect(slide, left, top, width, height, fill_color)
    card.line.color.rgb = border
    card.line.width = Pt(1.2)
    add_text(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.45), value, font_name="Aptos Display", font_size=24, color=accent, bold=True)
    add_text(slide, left + Inches(0.18), top + Inches(0.56), width - Inches(0.36), height - Inches(0.68), label, font_size=12, color=text_color)
    return card


def add_card_title(slide, left, top, width, title, text, accent=ORANGE):
    card = add_rect(slide, left, top, width, Inches(1.24), WHITE)
    card.line.color.rgb = RGBColor(226, 232, 240)
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.1), Inches(1.24))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.color.rgb = accent
    add_text(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.26), Inches(0.26), title, font_size=15, color=INK, bold=True)
    add_text(slide, left + Inches(0.18), top + Inches(0.42), width - Inches(0.26), Inches(0.66), text, font_size=11, color=SLATE)
    return card


def add_picture_cover(slide, image_path: Path, left, top, width, height):
    with Image.open(image_path) as im:
        img_w, img_h = im.size
    pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        new_width = img_h * box_ratio
        crop = (img_w - new_width) / (2 * img_w)
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        new_height = img_w / box_ratio
        crop = (img_h - new_height) / (2 * img_h)
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


def add_picture_contain(slide, image_path: Path, left, top, width, height, with_frame=True):
    if with_frame:
        frame = add_rect(slide, left, top, width, height, WHITE)
        frame.line.color.rgb = RGBColor(226, 232, 240)
    with Image.open(image_path) as im:
        img_w, img_h = im.size
    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        pic_w = width
        pic_h = int(width / img_ratio)
        pic_left = left
        pic_top = top + int((height - pic_h) / 2)
    else:
        pic_h = height
        pic_w = int(height * img_ratio)
        pic_top = top
        pic_left = left + int((width - pic_w) / 2)
    return slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_w, height=pic_h)


def add_path_caption(slide, left, top, width, text, dark=False):
    add_text(slide, left, top, width, Inches(0.16), text, font_name="Consolas", font_size=8, color=RGBColor(203, 213, 225) if dark else MUTED)


def add_code_box(slide, left, top, width, height, title, code):
    card = add_rect(slide, left, top, width, height, STEEL)
    card.line.color.rgb = ORANGE
    add_text(slide, left + Inches(0.14), top + Inches(0.08), width - Inches(0.28), Inches(0.24), title, font_size=12, color=WHITE, bold=True)
    add_text(slide, left + Inches(0.14), top + Inches(0.36), width - Inches(0.28), height - Inches(0.44), code, font_name="Consolas", font_size=10, color=RGBColor(226, 232, 240))
    return card


def pipeline_box(slide, left, top, text, fill_color):
    box = add_rect(slide, left, top, Inches(1.05), Inches(0.9), fill_color)
    box.line.color.rgb = fill_color
    add_text(slide, left + Inches(0.05), top + Inches(0.12), Inches(0.95), Inches(0.6), text, font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    return box


def flow_box(slide, left, top, width, height, title, text, color):
    box = add_rect(slide, left, top, width, height, WHITE)
    box.line.color.rgb = color
    box.line.width = Pt(1.2)
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.09), height)
    marker.fill.solid()
    marker.fill.fore_color.rgb = color
    marker.line.color.rgb = color
    add_text(slide, left + Inches(0.16), top + Inches(0.12), width - Inches(0.22), Inches(0.26), title, font_size=14, color=INK, bold=True)
    add_text(slide, left + Inches(0.16), top + Inches(0.42), width - Inches(0.22), height - Inches(0.5), text, font_size=11, color=SLATE)
    return box


def export_pdf_with_powerpoint(pptx_path: Path, pdf_path: Path) -> bool:
    script = rf"""
$ErrorActionPreference = 'Stop'
$app = New-Object -ComObject PowerPoint.Application
$app.Visible = -1
$presentation = $app.Presentations.Open('{pptx_path}', -1, -1, 0)
$presentation.SaveAs('{pdf_path}', 32)
$presentation.Close()
$app.Quit()
"""
    try:
        subprocess.run(["powershell", "-NoProfile", "-Command", script], check=True, capture_output=True, text=True)
        return True
    except subprocess.CalledProcessError:
        return False


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "PrediTeq - Soutenance jury"
    prs.core_properties.subject = "Presentation PrediTeq"
    prs.core_properties.author = "Codex"
    blank = prs.slide_layouts[6]

    holdout_rmse = METRICS["rul_regression"]["holdout"]["rmse_days"]
    holdout_r2 = METRICS["rul_regression"]["holdout"]["r2"]
    hybrid_f1 = METRICS["anomaly_detection"]["hybrid_ensemble"]["f1"]
    cv_r2 = RUL_CV["cross_validation_groupkfold"]["r2_mean"]
    cv_rmse = RUL_CV["cross_validation_groupkfold"]["rmse_days_mean"]
    cmapss_r2 = CMAPSS["r2"]
    cmapss_rmse = CMAPSS["rmse_cycles"]
    lead_mean_days = LEAD["lead_time_days_equiv_ui"]["mean_days"]
    n_trajectories = int(RUL_CV["config"]["n_trajectories"])
    motor_rpm = int(CFG.MOTOR_SPEED_RPM)
    load_max_kg = int(CFG.LOAD_MAX_KG)
    cycle_s = int(CFG.T_CYCLE_S)

    a1 = DEMO_MACHINE_SCENARIOS["ASC-A1"]
    b2 = DEMO_MACHINE_SCENARIOS["ASC-B2"]
    c3 = DEMO_MACHINE_SCENARIOS["ASC-C3"]

    total = 20

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_rect(slide, Inches(7.75), Inches(0), Inches(5.58), Inches(7.5), RGBColor(6, 11, 25), radius=False)
    add_picture_cover(slide, IMG / "photo_machine_aroteq.png", Inches(7.75), Inches(0), Inches(5.58), Inches(7.5))
    add_rect(slide, Inches(0.72), Inches(0.72), Inches(0.12), Inches(4.9), ORANGE, radius=False)
    add_text(slide, Inches(1.05), Inches(0.86), Inches(5.9), Inches(1.45), "PrediTeq", font_name="Aptos Display", font_size=30, color=WHITE, bold=True)
    add_text(slide, Inches(1.05), Inches(1.78), Inches(5.9), Inches(1.2), "Maintenance predictive intelligente pour un stockeur vertical rotatif", font_name="Aptos Display", font_size=22, color=RGBColor(226, 232, 240), bold=True)
    add_text(slide, Inches(1.05), Inches(3.02), Inches(5.9), Inches(0.82), "Projet PFE: Machine Learning + IoT + application web\nDu signal machine a la decision de maintenance", font_size=18, color=RGBColor(203, 213, 225))
    add_tag(slide, Inches(1.05), Inches(4.22), "Simulation", TEAL)
    add_tag(slide, Inches(2.42), Inches(4.22), "ML", ORANGE)
    add_tag(slide, Inches(3.79), Inches(4.22), "Runtime", SKY)
    add_tag(slide, Inches(5.16), Inches(4.22), "Web app", GOLD, text_color=INK)
    add_picture_contain(slide, IMG / "isamm.jpg", Inches(1.05), Inches(5.5), Inches(1.4), Inches(0.72), with_frame=False)
    add_picture_contain(slide, IMG / "aroteq-logo.png", Inches(2.62), Inches(5.52), Inches(1.82), Inches(0.68), with_frame=False)
    add_text(slide, Inches(1.05), Inches(6.34), Inches(5.8), Inches(0.5), "Soutenance jury | Version generee depuis les artefacts verifies du depot", font_size=11, color=RGBColor(203, 213, 225))
    add_footer(slide, 1, total)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Le probleme industriel", "Pourquoi PrediTeq existe et a quoi il sert")
    add_picture_cover(slide, IMG / "photo_machine_aroteq.png", Inches(0.75), Inches(1.45), Inches(4.0), Inches(4.9))
    add_bullets(
        slide,
        Inches(5.05),
        Inches(1.45),
        Inches(3.15),
        Inches(2.05),
        [
            "Intervenir trop tard peut mener a une panne ou a un arret.",
            "Intervenir trop tot cree un cout de maintenance inutile.",
            "Il faut donc surveiller l'etat de sante et anticiper le moment critique.",
        ],
        font_size=18,
    )
    add_card_title(slide, Inches(8.45), Inches(1.52), Inches(4.1), "Deux sorties cle", "Health Index = note de sante entre 0 et 1. RUL = temps restant avant la zone critique.", accent=ORANGE)
    add_card_title(slide, Inches(5.05), Inches(3.78), Inches(3.58), "Cas cible", "Stockeur vertical rotatif industriel AROTEQ utilise dans un contexte de production.", accent=TEAL)
    add_card_title(slide, Inches(8.8), Inches(3.78), Inches(3.75), "Valeur metier", "Aider a prioriser la maintenance, planifier et reduire les interventions mal synchronisees.", accent=SKY)
    add_picture_contain(slide, IMG / "schema_positionnement_stage_aroteq.png", Inches(5.05), Inches(5.15), Inches(7.5), Inches(1.1))
    add_footer(slide, 2, total)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Notre demarche methodologique", "Donnees limitees au depart, progression par phases, collecte reelle en parallele")
    add_card_title(
        slide,
        Inches(0.88),
        Inches(1.46),
        Inches(3.72),
        "1. Partir du reel",
        "Le projet part d'une vraie machine, mais sans longues historiques de pannes annotees au debut. Il fallait donc avancer sans attendre des annees.",
        accent=ORANGE,
    )
    add_card_title(
        slide,
        Inches(4.8),
        Inches(1.46),
        Inches(3.72),
        "2. Simuler de facon realiste",
        "Nous avons construit un socle d'entrainement a partir du moteur, du cycle, des charges, du bruit capteur et d'un contexte plausible.",
        accent=TEAL,
    )
    add_card_title(
        slide,
        Inches(8.72),
        Inches(1.46),
        Inches(3.72),
        "3. Avancer par phases",
        "Nous n'avons pas tout modele d'un coup : signaux, indicateurs, anomalies, note de sante, puis seulement temps restant.",
        accent=SKY,
    )
    add_card_title(
        slide,
        Inches(0.88),
        Inches(3.62),
        Inches(5.45),
        "4. Valider chaque bloc",
        "La chaine est verifiee sur un test separe, sur des trajectoires completes, sur NASA CMAPSS et par calibration. NASA sert ici de repere externe, pas de donnees principales du stockeur.",
        accent=ORANGE,
    )
    add_card_title(
        slide,
        Inches(6.92),
        Inches(3.62),
        Inches(5.52),
        "5. Collecter le reel en parallele",
        "Le flux LabVIEW / PLC vers le backend, l'historisation cloud et le journal maintenance preparent deja la phase suivante : raffiner et annoter avec du vrai terrain.",
        accent=TEAL,
    )
    add_rect(slide, Inches(0.88), Inches(5.85), Inches(11.56), Inches(0.62), STEEL)
    add_text(
        slide,
        Inches(1.08),
        Inches(5.98),
        Inches(11.15),
        Inches(0.26),
        "Le choix des variables vient aussi du terrain : plaque moteur, cycle reel et logique charge -> puissance -> courant -> echauffement discutee avec le technicien.",
        font_size=11,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 3, total)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Les donnees d'entrainement simulees", "Ce que contient vraiment le jeu cree a l'etape 1")
    add_bullets(
        slide,
        Inches(0.88),
        Inches(1.46),
        Inches(4.15),
        Inches(1.75),
        [
            "Une trajectoire = l'histoire complete d'une machine simulee dans le temps.",
            "Chaque trajectoire combine un profil de degradation et un cas de charge.",
            "Le script exporte ensuite les signaux utiles pour tout le pipeline ML.",
        ],
        font_size=16,
    )
    add_stat_card(slide, Inches(0.88), Inches(3.55), Inches(1.55), Inches(1.05), f"{n_trajectories}", "trajectoires", ORANGE)
    add_stat_card(slide, Inches(2.6), Inches(3.55), Inches(1.55), Inches(1.05), "4", "profils", TEAL)
    add_stat_card(slide, Inches(4.32), Inches(3.55), Inches(1.55), Inches(1.05), "20", "cas de charge", SKY)
    add_card_title(
        slide,
        Inches(0.88),
        Inches(4.95),
        Inches(4.95),
        "Colonnes cle",
        "trajectory_id, profile, load_kg, t_seconds, rms_mms, power_kw, current_a, temp_c, humidity_rh, phase, plus des metadonnees internes de simulation.",
        accent=ORANGE,
    )
    add_text(slide, Inches(0.9), Inches(6.06), Inches(4.9), Inches(0.34), "Sortie du script: prediteq_ml/data/raw/trajectories.csv", font_name="Consolas", font_size=10, color=SLATE)
    add_picture_contain(slide, IMG_FREE / "proof_code_07_training_dataset.png", Inches(6.12), Inches(1.44), Inches(6.2), Inches(4.95))
    add_path_caption(slide, Inches(6.14), Inches(6.42), Inches(5.8), "prediteq_ml/steps/step1_simulate.py")
    add_footer(slide, 4, total)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Pourquoi ces variables et comment elles donnent ces resultats", "Le choix des variables suit une chaine de cause a effet lisible")
    flow_box(slide, Inches(0.88), Inches(1.55), Inches(2.45), Inches(1.18), "Charge", "plus la charge est lourde, plus la montee demande d'effort", ORANGE)
    add_text(slide, Inches(3.4), Inches(1.98), Inches(0.25), Inches(0.25), ">", font_name="Aptos Display", font_size=22, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    flow_box(slide, Inches(3.68), Inches(1.55), Inches(2.45), Inches(1.18), "Effort electrique", "la puissance monte, donc le courant monte aussi", TEAL)
    add_text(slide, Inches(6.2), Inches(1.98), Inches(0.25), Inches(0.25), ">", font_name="Aptos Display", font_size=22, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    flow_box(slide, Inches(6.48), Inches(1.55), Inches(2.45), Inches(1.18), "Vieillissement", "l'echauffement et l'usure font baisser le Health Index", SKY)
    add_text(slide, Inches(9.0), Inches(1.98), Inches(0.25), Inches(0.25), ">", font_name="Aptos Display", font_size=22, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    flow_box(slide, Inches(9.28), Inches(1.55), Inches(3.12), Inches(1.18), "Signaux visibles", "vibration, puissance, courant, temperature et humidite evoluent", ORANGE)
    add_card_title(
        slide,
        Inches(0.88),
        Inches(3.2),
        Inches(4.25),
        "Pourquoi ces variables ?",
        "Nous avons garde des variables soit mesurables en live, soit directement derivees du comportement machine. Ainsi, l'entrainement offline et le runtime utilisent le meme langage physique.",
        accent=TEAL,
    )
    add_card_title(
        slide,
        Inches(0.88),
        Inches(4.7),
        Inches(4.25),
        "Message non technique",
        "Nous n'avons pas choisi des colonnes arbitraires. Nous avons choisi les mesures qui changent quand la machine travaille plus fort ou se degrade.",
        accent=ORANGE,
    )
    add_picture_contain(slide, IMG_FREE / "proof_code_01_sim_constants.png", Inches(5.45), Inches(3.08), Inches(3.3), Inches(3.1))
    add_picture_contain(slide, IMG_FREE / "proof_code_02_sim_power_current.png", Inches(8.95), Inches(3.08), Inches(3.3), Inches(3.1))
    add_path_caption(slide, Inches(5.48), Inches(6.24), Inches(3.2), "prediteq_ml/config.py")
    add_path_caption(slide, Inches(8.98), Inches(6.24), Inches(3.2), "prediteq_ml/steps/step1_simulate.py")
    add_footer(slide, 5, total)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "La logique des 3 machines simulees", "Meme base projet, mais 3 contextes d'usage pour montrer 3 etats de sante")
    a1_low, a1_high = a1["load_band_kg"]
    b2_low, b2_high = b2["load_band_kg"]
    c3_low, c3_high = c3["load_band_kg"]
    add_card_title(
        slide,
        Inches(0.88),
        Inches(1.45),
        Inches(4.35),
        "ASC-A1 - machine protegee",
        f"{a1['site']} | profil {a1['profile']} | charge {a1_low}-{a1_high} kg | cible HI {a1['target_hi']:.2f}. Usage leger, environnement plus sec, peu de surcharge.",
        accent=TEAL,
    )
    add_card_title(
        slide,
        Inches(0.88),
        Inches(2.92),
        Inches(4.35),
        "ASC-B2 - machine sous surveillance",
        f"{b2['site']} | profil {b2['profile']} | charge {b2_low}-{b2_high} kg | cible HI {b2['target_hi']:.2f}. Trafic mixte, demi-charges frequentes, stress moyen.",
        accent=ORANGE,
    )
    add_card_title(
        slide,
        Inches(0.88),
        Inches(4.39),
        Inches(4.35),
        "ASC-C3 - machine critique",
        f"{c3['site']} | profil {c3['profile']} | charge {c3_low}-{c3_high} kg | cible HI {c3['target_hi']:.2f}. Ligne intensive, charges lourdes, usure et stress eleves.",
        accent=RED,
    )
    add_text(slide, Inches(0.92), Inches(5.94), Inches(4.25), Inches(0.36), "Message simple: ce ne sont pas 3 graines aleatoires. Le simulateur change vraiment les charges, le stress, l'usure et le point de depart.", font_size=11, color=SLATE)
    add_picture_contain(slide, IMG_FREE / "proof_code_08_demo_machine_scenarios.png", Inches(5.55), Inches(1.45), Inches(6.75), Inches(2.2))
    add_picture_contain(slide, IMG_FREE / "proof_code_09_demo_machine_runtime_logic.png", Inches(5.55), Inches(3.95), Inches(6.75), Inches(2.2))
    add_path_caption(slide, Inches(5.58), Inches(3.72), Inches(4.8), "prediteq_api/demo_scenarios.py")
    add_path_caption(slide, Inches(5.58), Inches(6.24), Inches(5.5), "prediteq_api/routers/simulator.py")
    add_footer(slide, 6, total)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_title_block(slide, "Pipeline ML en langage simple", "Chaque bloc a un role clair, puis le backend charge les artefacts exportes", dark=True)
    steps = [
        ("1\nSimuler", TEAL),
        ("2\nPreparer", SKY),
        ("3\nDetecter", ORANGE),
        ("4\nNoter la\nsante", TEAL),
        ("5\nEstimer le\ntemps", SKY),
        ("6\nMesurer", ORANGE),
        ("6B\nTester\nNASA", TEAL),
        ("6C\nCalibrer", SKY),
        ("7\nExporter", ORANGE),
    ]
    x = 0.76
    for idx, (label, color) in enumerate(steps):
        pipeline_box(slide, Inches(x), Inches(2.2), label, color)
        if idx < len(steps) - 1:
            add_text(slide, Inches(x + 1.11), Inches(2.46), Inches(0.22), Inches(0.2), ">", font_name="Aptos Display", font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += 1.35
    add_rect(slide, Inches(0.88), Inches(4.05), Inches(11.58), Inches(1.55), STEEL)
    add_text(slide, Inches(1.1), Inches(4.28), Inches(11.1), Inches(0.4), "Lecture simple pour le jury", font_name="Aptos Display", font_size=18, color=WHITE, bold=True)
    add_bullets(
        slide,
        Inches(1.1),
        Inches(4.72),
        Inches(11.0),
        Inches(0.7),
        [
            "Nous avons avance par blocs: simulation -> indicateurs -> anomalies -> HI -> RUL.",
            "Chaque bloc a ete valide avant le suivant, puis exporte vers le runtime.",
            "Le backend charge ces exports; il ne reentraine pas en live.",
        ],
        font_size=16,
        color=RGBColor(226, 232, 240),
        bullet_color=ORANGE,
    )
    add_text(slide, Inches(0.9), Inches(6.1), Inches(11.6), Inches(0.35), "Ordre a citer: step1 -> step2 -> step3 -> step4 -> step5 -> step6 -> step6b -> step6c -> step7", font_name="Consolas", font_size=11, color=RGBColor(203, 213, 225))
    add_footer(slide, 7, total)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Des mesures brutes vers une note de sante", "Une feature = un indicateur calcule a partir des mesures brutes")
    add_bullets(
        slide,
        Inches(0.85),
        Inches(1.46),
        Inches(3.3),
        Inches(2.15),
        [
            "Au depart, on a surtout vibration, puissance, temperature et humidite.",
            "Le pipeline calcule des indicateurs intermediaires, par exemple une moyenne, une variation ou une duree de montee.",
            "Ces indicateurs servent a reconnaitre un comportement inhabituel.",
            "Le resultat final devient une note de sante HI plus facile a lire.",
        ],
        font_size=15,
    )
    add_picture_contain(slide, IMG_FREE / "proof_code_03_ml_features.png", Inches(4.32), Inches(1.44), Inches(3.95), Inches(1.9))
    add_picture_contain(slide, IMG_FREE / "proof_code_04_ml_hybrid_score.png", Inches(8.47), Inches(1.44), Inches(3.95), Inches(1.9))
    add_path_caption(slide, Inches(4.35), Inches(3.38), Inches(3.8), "prediteq_ml/steps/step2_preprocess.py")
    add_path_caption(slide, Inches(8.5), Inches(3.38), Inches(3.8), "prediteq_ml/steps/step3_isolation_forest.py")
    add_picture_contain(slide, IMG / "plot3_anomaly_timeline.png", Inches(0.85), Inches(3.72), Inches(11.57), Inches(2.35))
    add_footer(slide, 8, total)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "RUL: estimer le temps restant", "RUL = temps restant avant qu'une action devienne necessaire")
    add_picture_contain(slide, IMG / "plot2_rul_scatter.png", Inches(0.85), Inches(1.45), Inches(6.35), Inches(4.75))
    add_bullets(
        slide,
        Inches(7.45),
        Inches(1.46),
        Inches(4.15),
        Inches(1.65),
        [
            "Le RUL repond a une question simple: combien de temps reste-t-il avant qu'il faille agir ?",
            "Il est construit a partir du Health Index observe par le systeme.",
            "Le franchissement persistant evite les faux changements dus au bruit.",
            "Le but est d'aider a planifier, pas seulement a declencher une alerte.",
        ],
        font_size=14,
    )
    add_stat_card(slide, Inches(7.5), Inches(3.36), Inches(1.95), Inches(1.08), f"{holdout_rmse:.2f} j", "erreur moy. sur test separe", ORANGE)
    add_stat_card(slide, Inches(9.62), Inches(3.36), Inches(1.95), Inches(1.08), f"{cv_r2:.3f}", "stabilite par trajectoires", TEAL)
    add_text(slide, Inches(7.5), Inches(4.62), Inches(4.1), Inches(0.4), f"Lead time moyen avant critique: {lead_mean_days:.1f} jours eq UI.\nVous n'avez pas besoin de retenir les sigles.", font_size=11, color=SLATE)
    add_picture_contain(slide, IMG_FREE / "proof_code_05_ml_rul_target.png", Inches(7.45), Inches(4.92), Inches(2.0), Inches(1.45))
    add_picture_contain(slide, IMG_FREE / "proof_code_06_ml_groupkfold.png", Inches(9.58), Inches(4.92), Inches(2.0), Inches(1.45))
    add_footer(slide, 9, total)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_title_block(slide, "Comment nous avons verifie les resultats", "Plusieurs controles pour eviter l'effet 'ca marche seulement chez nous'", dark=True)
    add_stat_card(slide, Inches(0.88), Inches(1.46), Inches(2.35), Inches(1.12), f"{hybrid_f1:.3f}", "qualite detection (F1)", ORANGE, dark=True)
    add_stat_card(slide, Inches(3.5), Inches(1.46), Inches(2.35), Inches(1.12), f"{holdout_r2:.3f}", "qualite prediction RUL (R2)", TEAL, dark=True)
    add_stat_card(slide, Inches(6.12), Inches(1.46), Inches(2.35), Inches(1.12), f"{cv_rmse:.2f} j", "erreur moyenne validation", SKY, dark=True)
    add_stat_card(slide, Inches(8.74), Inches(1.46), Inches(2.75), Inches(1.12), f"{cmapss_r2:.3f} | {cmapss_rmse:.1f}", "test externe NASA", GOLD, dark=True)
    add_picture_contain(slide, IMG / "plot6_cmapss.png", Inches(0.88), Inches(3.02), Inches(5.85), Inches(3.0), with_frame=False)
    add_picture_contain(slide, IMG / "plot7_calibration.png", Inches(6.95), Inches(3.02), Inches(5.45), Inches(3.0), with_frame=False)
    add_text(slide, Inches(0.92), Inches(6.15), Inches(11.5), Inches(0.38), "Pas besoin de retenir tous les sigles: l'idee est que la chaine reste stable sur un test separe, sur des trajectoires entieres et sur un benchmark externe NASA.", font_size=11, color=RGBColor(203, 213, 225))
    add_footer(slide, 10, total)

    # Slide 11
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Du signal a l'application", "Que se passe-t-il entre une mesure et l'ecran du technicien ?")
    flow_box(slide, Inches(0.88), Inches(1.75), Inches(2.1), Inches(1.35), "Sources", "capteurs, simulateur ou CSV LabVIEW de demo", ORANGE)
    add_text(slide, Inches(3.05), Inches(2.22), Inches(0.3), Inches(0.3), ">", font_name="Aptos Display", font_size=24, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    flow_box(slide, Inches(3.38), Inches(1.75), Inches(1.85), Inches(1.35), "Bridge", "remet les champs\nau bon format", TEAL)
    add_text(slide, Inches(5.29), Inches(2.22), Inches(0.3), Inches(0.3), ">", font_name="Aptos Display", font_size=24, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    flow_box(slide, Inches(5.62), Inches(1.75), Inches(2.1), Inches(1.35), "FastAPI runtime", "charge les modeles exportes\net calcule note de sante, zone et temps restant", SKY)
    add_text(slide, Inches(7.79), Inches(2.22), Inches(0.3), Inches(0.3), ">", font_name="Aptos Display", font_size=24, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    flow_box(slide, Inches(8.12), Inches(1.75), Inches(1.9), Inches(1.35), "Supabase", "garde les seuils live\net le contexte machine", ORANGE)
    add_text(slide, Inches(10.1), Inches(2.22), Inches(0.3), Inches(0.3), ">", font_name="Aptos Display", font_size=24, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    flow_box(slide, Inches(10.45), Inches(1.75), Inches(1.95), Inches(1.35), "App web", "affiche et aide\na decider", TEAL)
    add_card_title(slide, Inches(0.88), Inches(3.55), Inches(3.65), "Point cle", "Le backend n'entraine rien. Il applique les modeles deja exportes depuis prediteq_ml/models.", accent=SKY)
    add_card_title(slide, Inches(4.75), Inches(3.55), Inches(3.65), "Pourquoi c'est rassurant", "La meme logique peut etre relue dans les fichiers ML offline puis retrouvee dans le runtime.", accent=ORANGE)
    add_picture_contain(slide, IMG / "technologies_aroteq.png", Inches(8.7), Inches(3.34), Inches(3.7), Inches(2.6))
    add_footer(slide, 11, total)

    # Slide 12
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Ce que voit l'utilisateur", "Le produit final transforme les calculs en decisions lisibles")
    add_picture_contain(slide, IMG / "capture_dashboard_prediteq.png", Inches(5.0), Inches(1.45), Inches(7.35), Inches(4.95))
    add_card_title(slide, Inches(0.88), Inches(1.6), Inches(3.65), "Dashboard", "Vue globale de la machine, des alertes et du contexte de fonctionnement.", accent=ORANGE)
    add_card_title(slide, Inches(0.88), Inches(3.0), Inches(3.65), "Diagnostics / RUL", "Voir la note de sante, la zone courante et le temps restant avant criticite.", accent=TEAL)
    add_card_title(slide, Inches(0.88), Inches(4.4), Inches(3.65), "Planner / Maintenance / Costs", "Prioriser l'action, planifier et comprendre l'impact economique.", accent=SKY)
    add_text(slide, Inches(0.9), Inches(6.12), Inches(3.8), Inches(0.3), "En demo, ouvrir: Dashboard -> Diagnostics -> Planner -> Maintenance / Costs", font_size=11, color=SLATE)
    add_footer(slide, 12, total)

    # Slide 13
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Integration machine reelle / PLC / LabVIEW", "ARO-01 suit deja le vrai chemin live; la source finale remplacera surtout le bloc d'acquisition")
    flow_box(slide, Inches(0.88), Inches(1.55), Inches(2.1), Inches(1.22), "LabVIEW / PLC", "acquisition des mesures sur la machine reelle", ORANGE)
    add_text(slide, Inches(3.05), Inches(1.98), Inches(0.3), Inches(0.3), ">", font_name="Aptos Display", font_size=24, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    flow_box(slide, Inches(3.38), Inches(1.55), Inches(1.95), Inches(1.22), "PC relais site", "ecrit un JSON ou un CSV, puis envoie", TEAL)
    add_text(slide, Inches(5.42), Inches(1.98), Inches(0.3), Inches(0.3), ">", font_name="Aptos Display", font_size=24, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    flow_box(slide, Inches(5.72), Inches(1.55), Inches(1.8), Inches(1.22), "MQTT / HTTP", "transport vers PrediTeq", SKY)
    add_text(slide, Inches(7.61), Inches(1.98), Inches(0.3), Inches(0.3), ">", font_name="Aptos Display", font_size=24, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    flow_box(slide, Inches(7.92), Inches(1.55), Inches(2.2), Inches(1.22), "FastAPI /ingest/live", "verifie la machine et injecte le payload", ORANGE)
    add_text(slide, Inches(10.2), Inches(1.98), Inches(0.3), Inches(0.3), ">", font_name="Aptos Display", font_size=24, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    flow_box(slide, Inches(10.55), Inches(1.55), Inches(1.85), Inches(1.22), "Frontend", "affiche HI, zone, alertes et RUL", TEAL)
    add_card_title(
        slide,
        Inches(0.88),
        Inches(3.15),
        Inches(3.9),
        "4 mesures minimales",
        "rms_mms, power_kw, temp_c et humidity_rh. current_a, load_kg et status restent utiles mais optionnels.",
        accent=ORANGE,
    )
    add_card_title(
        slide,
        Inches(0.88),
        Inches(4.68),
        Inches(3.9),
        "Ce qui change plus tard",
        "Aujourd'hui, la soutenance utilise une source CSV LabVIEW de demo. Demain, la vraie acquisition PLC/LabVIEW remplacera seulement cette source.",
        accent=TEAL,
    )
    add_text(slide, Inches(5.08), Inches(6.42), Inches(6.9), Inches(0.22), "ARO-01 n'est pas une 4e machine simulee : meme moteur runtime, sans overrides simulateur.", font_size=11, color=SLATE)
    add_picture_contain(slide, IMG_FREE / "proof_code_10_live_ingest_bridge.png", Inches(5.05), Inches(3.0), Inches(7.3), Inches(3.1))
    add_path_caption(slide, Inches(5.08), Inches(6.16), Inches(6.6), "prediteq_api/routers/live_ingest.py + prediteq_api/RELAY_PC_SETUP_SIMPLE.md")
    add_footer(slide, 13, total)

    # Slide 14
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_title_block(slide, "Demonstration live facile a suivre", "ARO-01 est d'abord preparee, puis le flux CSV live continue sur la meme machine", dark=True)
    add_text(slide, Inches(0.88), Inches(1.35), Inches(11.6), Inches(0.25), "bootstrap ARO-01 -> CSV LabVIEW de demonstration -> mqtt_bridge_sender.py -> FastAPI -> application web", font_name="Consolas", font_size=12, color=RGBColor(203, 213, 225))
    add_picture_contain(slide, IMG / "validation_pc2_sender_clean.png", Inches(0.88), Inches(1.82), Inches(3.85), Inches(3.88), with_frame=False)
    add_picture_contain(slide, IMG / "validation_pc1_backend_clean.png", Inches(4.75), Inches(1.82), Inches(3.85), Inches(3.88), with_frame=False)
    add_picture_contain(slide, IMG / "validation_pc1_app_aro01_clean.png", Inches(8.62), Inches(1.82), Inches(3.83), Inches(3.88), with_frame=False)
    add_text(slide, Inches(1.0), Inches(5.85), Inches(3.6), Inches(0.32), "1. Source de donnees", font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(4.9), Inches(5.85), Inches(3.5), Inches(0.32), "2. Calcul backend", font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(8.75), Inches(5.85), Inches(3.5), Inches(0.32), "3. Resultat dans l'app", font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 14, total)

    # Slide 15
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Conclusion et prochaines etapes", "Ce qui est deja solide, puis ce qui doit encore etre rapproche du terrain")
    add_rect(slide, Inches(0.88), Inches(1.45), Inches(5.55), Inches(4.85), WHITE)
    add_rect(slide, Inches(6.85), Inches(1.45), Inches(5.55), Inches(4.85), WHITE)
    add_text(slide, Inches(1.1), Inches(1.7), Inches(2.9), Inches(0.3), "Deja en place", font_name="Aptos Display", font_size=18, color=ORANGE, bold=True)
    add_bullets(
        slide,
        Inches(1.08),
        Inches(2.05),
        Inches(5.0),
        Inches(3.9),
        [
            "Jeu d'entrainement simule relie a la logique machine.",
            "Trois machines demo distinctes pour raconter trois regimes d'usage.",
            "Pipeline ML ordonne et mesure sur plusieurs tests.",
            "Backend runtime + application web + integration live deja structuree.",
        ],
        font_size=16,
    )
    add_text(slide, Inches(7.07), Inches(1.7), Inches(3.2), Inches(0.3), "Suite logique", font_name="Aptos Display", font_size=18, color=TEAL, bold=True)
    add_bullets(
        slide,
        Inches(7.05),
        Inches(2.05),
        Inches(5.0),
        Inches(3.9),
        [
            "Augmenter le volume de donnees terrain reelles et les annotations.",
            "Remplacer progressivement la source CSV demo par une acquisition industrielle directe.",
            "Affiner encore la calibration et les seuils par contexte machine.",
            "Etendre la generalisation vers d'autres equipements et profils d'usage.",
        ],
        font_size=16,
    )
    add_rect(slide, Inches(0.88), Inches(6.45), Inches(11.52), Inches(0.45), STEEL)
    add_text(slide, Inches(1.1), Inches(6.52), Inches(11.1), Inches(0.22), "Phrase de fin: PrediTeq relie de facon coherente la physique machine, le Machine Learning, le runtime IoT et la decision dans une application web.", font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_footer(slide, 15, total)

    # Slide 16
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Backup - jeu d'entrainement et preuves simulation", "step1_simulate.py + config.py: ce que nous presentons vient du vrai code")
    add_picture_contain(slide, IMG_FREE / "proof_code_07_training_dataset.png", Inches(0.88), Inches(1.45), Inches(5.55), Inches(4.85))
    add_picture_contain(slide, IMG_FREE / "proof_code_01_sim_constants.png", Inches(6.72), Inches(1.45), Inches(5.55), Inches(2.2))
    add_picture_contain(slide, IMG_FREE / "proof_code_02_sim_power_current.png", Inches(6.72), Inches(4.12), Inches(5.55), Inches(2.2))
    add_path_caption(slide, Inches(0.9), Inches(6.38), Inches(4.8), "prediteq_ml/steps/step1_simulate.py")
    add_path_caption(slide, Inches(6.75), Inches(3.74), Inches(4.8), "prediteq_ml/config.py")
    add_path_caption(slide, Inches(6.75), Inches(6.38), Inches(5.0), "prediteq_ml/steps/step1_simulate.py")
    add_footer(slide, 16, total)

    # Slide 17
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Backup - preuves des 3 machines de demo", "demo_scenarios.py + simulator.py: charges, stress, usure et cibles HI")
    add_picture_contain(slide, IMG_FREE / "proof_code_08_demo_machine_scenarios.png", Inches(0.88), Inches(1.45), Inches(11.45), Inches(2.15))
    add_picture_contain(slide, IMG_FREE / "proof_code_09_demo_machine_runtime_logic.png", Inches(0.88), Inches(3.95), Inches(11.45), Inches(2.15))
    add_path_caption(slide, Inches(0.92), Inches(3.66), Inches(5.4), "prediteq_api/demo_scenarios.py")
    add_path_caption(slide, Inches(0.92), Inches(6.16), Inches(5.7), "prediteq_api/routers/simulator.py")
    add_footer(slide, 17, total)

    # Slide 18
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Backup - preuves ML actuelles", "step2, step3, step4 et step5: les vraies briques de la chaine ML")
    add_picture_contain(slide, IMG_FREE / "proof_code_03_ml_features.png", Inches(0.88), Inches(1.45), Inches(5.4), Inches(2.25))
    add_picture_contain(slide, IMG_FREE / "proof_code_04_ml_hybrid_score.png", Inches(6.95), Inches(1.45), Inches(5.4), Inches(2.25))
    add_picture_contain(slide, IMG_FREE / "proof_code_05_ml_rul_target.png", Inches(0.88), Inches(3.95), Inches(5.4), Inches(2.25))
    add_picture_contain(slide, IMG_FREE / "proof_code_06_ml_groupkfold.png", Inches(6.95), Inches(3.95), Inches(5.4), Inches(2.25))
    add_footer(slide, 18, total)

    # Slide 19
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_title_block(slide, "Backup - integration live et commandes", "Contrat HTTP/MQTT, puis commandes minimales pour la demo locale", dark=True)
    add_picture_contain(slide, IMG_FREE / "proof_code_10_live_ingest_bridge.png", Inches(0.88), Inches(1.45), Inches(5.35), Inches(4.85))
    add_code_box(slide, Inches(6.45), Inches(1.45), Inches(5.9), Inches(1.72), "Backend + frontend", "cd prediteq_api\npip install -r requirements.txt\nuvicorn main:app --reload\n\ncd prediteq_frontend\nnpm install\nnpm run dev\n# http://127.0.0.1:8080")
    add_code_box(slide, Inches(6.45), Inches(3.55), Inches(5.9), Inches(2.72), "Demo LabVIEW / PC relais", r"cd prediteq_api\npython scripts/generate_labview_demo_csv.py --machine-id ARO-01 --scenario surveillance\npython scripts/setup_real_machine_demo.py --machine-id ARO-01 --name 'Machine reelle' --scenario surveillance\npython scripts/replay_labview_demo_csv.py --input scripts/sample_data/ARO-01_labview_demo_template.csv --output C:\labview\prediteq_log.csv --interval 1.0\npython scripts/mqtt_bridge_sender.py --transport mqtt --mode csv-last-row --machine-id ARO-01 --csv-path C:\labview\prediteq_log.csv")
    add_footer(slide, 19, total)

    # Slide 20
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_title_block(slide, "Backup - fichiers source de verite", "Ou montrer le code actuel en moins de 30 secondes")
    add_code_box(slide, Inches(0.88), Inches(1.45), Inches(3.75), Inches(1.95), "Jeu d'entrainement", "prediteq_ml/config.py\nprediteq_ml/steps/step1_simulate.py")
    add_code_box(slide, Inches(4.78), Inches(1.45), Inches(3.75), Inches(1.95), "3 machines demo", "prediteq_api/demo_scenarios.py\nprediteq_api/routers/simulator.py")
    add_code_box(slide, Inches(8.68), Inches(1.45), Inches(3.67), Inches(1.95), "ML / resultats", "prediteq_ml/steps/step2_preprocess.py\nprediteq_ml/steps/step3_isolation_forest.py\nprediteq_ml/steps/step4_health_index.py\nprediteq_ml/steps/step5_rul_model.py")
    add_code_box(slide, Inches(1.15), Inches(4.0), Inches(4.2), Inches(1.95), "Runtime backend", "prediteq_api/ml/engine_manager.py\nprediteq_api/routers/diagnostics_rul.py\nprediteq_api/routers/live_ingest.py")
    add_code_box(slide, Inches(7.0), Inches(4.0), Inches(4.2), Inches(1.95), "Flux live / PC relais", "prediteq_api/scripts/setup_real_machine_demo.py\nprediteq_api/scripts/generate_labview_demo_csv.py\nprediteq_api/scripts/mqtt_bridge_sender.py\nprediteq_api/LABVIEW_CSV_BRIDGE_DEMO.md")
    add_footer(slide, 20, total)

    return prs


def main() -> int:
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(str(OUTPUT_PPTX))
    exported_pdf = export_pdf_with_powerpoint(OUTPUT_PPTX, OUTPUT_PDF)
    print(f"PPTX_CREATED={OUTPUT_PPTX}")
    print(f"PDF_CREATED={OUTPUT_PDF if exported_pdf else 'NO'}")
    print(f"SLIDE_COUNT={len(prs.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
